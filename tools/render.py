#!/usr/bin/env python3
"""
glosa/tools/render.py — document-rendering layer for the glosa toolchain.

Readout discipline: this tool renders documents. It does not change, verify,
or upgrade the epistemic tier of any claim inside them. Every output carries
a footer/closing line stating the tier is unchanged (Dr unless the source
markdown states otherwise) and naming the render.py version that produced it.

Stdlib + `markdown` + `python-docx` + `python-pptx` only for the direct
routes; LibreOffice (`soffice --headless`) is used as the HTML/DOCX -> PDF
and the DOCX fallback engine. No shell strings are ever built for
subprocess — every external call is an argv list.

Commands:
    render.py md2html   in.md out.html [--lang th|en] [--css]
    render.py md2pdf    in.md out.pdf
    render.py md2docx   in.md out.docx
    render.py md2pptx   in.md out.pptx
    render.py tex2pdf   main.tex
    render.py all       paper/main_en.md paper/main_th.md --out dist/
    render.py check
    render.py deliver-manifest out.json [--dir dist]
"""
from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

VERSION = "0.1.0"
FOOTER_TEXT = f"Rendered with glosa render.py {VERSION} — tier of content unchanged (Dr unless stated)"


# --------------------------------------------------------------------------
# Front matter + markdown parsing helpers
# --------------------------------------------------------------------------

def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def split_front_matter(text: str) -> tuple[dict, str]:
    """Very small YAML-ish front matter parser: `key: value` lines between
    a leading `---` and a closing `---`. Not a full YAML parser — the glosa
    paper front matter is flat key:value pairs only."""
    meta: dict[str, str] = {}
    if text.startswith("---"):
        parts = text.split("\n---", 1)
        if len(parts) == 2:
            head = parts[0][3:]
            body = parts[1]
            if body.startswith("\n"):
                body = body[1:]
            for line in head.splitlines():
                line = line.strip()
                if not line or ":" not in line:
                    continue
                k, _, v = line.partition(":")
                meta[k.strip()] = v.strip().strip('"').strip("'")
            return meta, body
    return meta, text


def has_thai(text: str) -> bool:
    return bool(re.search(r"[\u0E00-\u0E7F]", text))


# --------------------------------------------------------------------------
# md2html
# --------------------------------------------------------------------------

CSS_TEMPLATE = """
@charset "UTF-8";
body {{
  font-family: {body_font};
  line-height: 1.6;
  max-width: 880px;
  margin: 2.5rem auto;
  padding: 0 1.5rem;
  color: #1a1a1a;
}}
h1, h2, h3, h4 {{ font-family: {head_font}; line-height: 1.3; }}
h1 {{ font-size: 1.9rem; border-bottom: 2px solid #333; padding-bottom: .3rem; }}
h2 {{ font-size: 1.4rem; margin-top: 2rem; }}
h3 {{ font-size: 1.15rem; }}
table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; }}
th, td {{ border: 1px solid #999; padding: .4rem .6rem; text-align: left; }}
th {{ background: #eee; }}
code, pre {{ font-family: "DejaVu Sans Mono", monospace; }}
pre {{ background: #f4f4f4; padding: .8rem; overflow-x: auto; border-radius: 4px; }}
blockquote {{ border-left: 4px solid #999; margin: 1rem 0; padding: .2rem 1rem; color: #444; }}
.glosa-footer {{ margin-top: 3rem; padding-top: .8rem; border-top: 1px solid #ccc;
  font-size: .8rem; color: #666; }}
"""

FONT_STACKS = {
    "th": {
        "body": '"Noto Sans Thai", "Noto Looped Thai", "Noto Serif Thai", sans-serif',
        "head": '"Noto Serif Thai", "Noto Sans Thai", serif',
    },
    "en": {
        "body": '"Noto Serif", Georgia, "Times New Roman", serif',
        "head": '"Noto Serif", Georgia, serif',
    },
}


def build_css(lang: str) -> str:
    fonts = FONT_STACKS.get(lang, FONT_STACKS["en"])
    return CSS_TEMPLATE.format(body_font=fonts["body"], head_font=fonts["head"])


def md_to_html_fragment(md_text: str) -> str:
    import markdown

    extensions = ["tables", "fenced_code", "sane_lists", "toc"]
    try:
        import markdown.extensions.footnotes  # noqa: F401
        extensions.append("footnotes")
    except Exception:
        pass
    return markdown.markdown(md_text, extensions=extensions)


def cmd_md2html(args: argparse.Namespace) -> int:
    in_path = Path(args.input)
    out_path = Path(args.output)
    text = read_text(in_path)
    meta, body = split_front_matter(text)

    lang = args.lang or meta.get("lang", "en")
    title = meta.get("title", in_path.stem)

    fragment = md_to_html_fragment(body)
    css = build_css(lang) if args.css or True else ""  # CSS always embedded per spec

    html = f"""<!doctype html>
<html lang="{lang}">
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>{css}</style>
</head>
<body>
{fragment}
<div class="glosa-footer">{FOOTER_TEXT}</div>
</body>
</html>
"""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(html, encoding="utf-8")
    print(f"wrote {out_path}")
    return 0


# --------------------------------------------------------------------------
# soffice helpers
# --------------------------------------------------------------------------

def soffice_convert(src: Path, out_dir: Path, target_filter: str) -> Path:
    """Run `soffice --headless --convert-to <target_filter> --outdir <out_dir> <src>`.
    Returns the path soffice produced (name derived from src stem)."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("soffice/libreoffice not found on PATH")
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [soffice, "--headless", "--convert-to", target_filter, "--outdir", str(out_dir), str(src)]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
    if result.returncode != 0:
        raise RuntimeError(f"soffice failed ({result.returncode}): {result.stdout}\n{result.stderr}")
    ext = target_filter.split(":", 1)[0]
    produced = out_dir / f"{src.stem}.{ext}"
    if not produced.exists():
        raise RuntimeError(f"soffice reported success but {produced} not found. stdout={result.stdout}")
    return produced


def cmd_md2pdf(args: argparse.Namespace) -> int:
    in_path = Path(args.input)
    out_path = Path(args.output)
    with tempfile.TemporaryDirectory() as td:
        tmp_html = Path(td) / (in_path.stem + ".html")
        html_args = argparse.Namespace(input=str(in_path), output=str(tmp_html), lang=None, css=True)
        cmd_md2html(html_args)
        produced = soffice_convert(tmp_html, Path(td), "pdf")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy(produced, out_path)
    print(f"wrote {out_path}")
    return 0


# --------------------------------------------------------------------------
# md2docx — direct python-docx build, soffice fallback
# --------------------------------------------------------------------------

def parse_markdown_blocks(body: str):
    """Yield simple block tuples for direct docx/pptx building:
    ('h', level, text) | ('p', text) | ('ul', [items]) | ('ol', [items]) |
    ('table', [[cells],...]) | ('code', text) | ('quote', text)
    Deliberately simple — glosa papers use headings/paragraphs/lists/tables/
    code fences/blockquotes; anything fancier degrades to a paragraph."""
    lines = body.splitlines()
    blocks = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        m = re.match(r"^(#{1,4})\s+(.*)$", stripped)
        if m:
            blocks.append(("h", len(m.group(1)), m.group(2).strip()))
            i += 1
            continue
        if stripped.startswith("```"):
            code_lines = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            blocks.append(("code", "\n".join(code_lines)))
            continue
        if stripped.startswith(">"):
            quote_lines = []
            while i < n and lines[i].strip().startswith(">"):
                quote_lines.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            blocks.append(("quote", "\n".join(quote_lines)))
            continue
        if re.match(r"^\|.*\|$", stripped):
            table_lines = []
            while i < n and re.match(r"^\|.*\|$", lines[i].strip()):
                table_lines.append(lines[i].strip())
                i += 1
            rows = []
            for tl in table_lines:
                if re.match(r"^\|[\s:\-|]+\|$", tl):
                    continue  # separator row
                cells = [c.strip() for c in tl.strip("|").split("|")]
                rows.append(cells)
            blocks.append(("table", rows))
            continue
        if re.match(r"^[-*]\s+", stripped):
            items = []
            while i < n and re.match(r"^[-*]\s+", lines[i].strip()):
                items.append(re.sub(r"^[-*]\s+", "", lines[i].strip()))
                i += 1
            blocks.append(("ul", items))
            continue
        if re.match(r"^\d+\.\s+", stripped):
            items = []
            while i < n and re.match(r"^\d+\.\s+", lines[i].strip()):
                items.append(re.sub(r"^\d+\.\s+", "", lines[i].strip()))
                i += 1
            blocks.append(("ol", items))
            continue
        # paragraph: gather until blank line
        para_lines = [stripped]
        i += 1
        while i < n and lines[i].strip() and not re.match(r"^(#{1,4})\s+|^```|^\||^[-*]\s+|^\d+\.\s+|^>", lines[i].strip()):
            para_lines.append(lines[i].strip())
            i += 1
        blocks.append(("p", " ".join(para_lines)))
    return blocks


def strip_inline_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1 (\2)", text)
    return text


def build_docx_direct(meta: dict, blocks, out_path: Path) -> None:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    title = meta.get("title")
    if title:
        doc.add_heading(strip_inline_md(title), level=0)
        subtitle_bits = [v for k, v in meta.items() if k in ("author", "date", "lang") and v]
        if subtitle_bits:
            doc.add_paragraph(" · ".join(subtitle_bits))

    for block in blocks:
        kind = block[0]
        if kind == "h":
            _, level, text = block
            doc.add_heading(strip_inline_md(text), level=min(level, 4))
        elif kind == "p":
            doc.add_paragraph(strip_inline_md(block[1]))
        elif kind == "quote":
            p = doc.add_paragraph(strip_inline_md(block[1]))
            p.style = doc.styles["Intense Quote"] if "Intense Quote" in [s.name for s in doc.styles] else p.style
        elif kind == "code":
            p = doc.add_paragraph(block[1])
            for run in p.runs:
                run.font.name = "DejaVu Sans Mono"
                run.font.size = Pt(9)
        elif kind == "ul":
            for item in block[1]:
                doc.add_paragraph(strip_inline_md(item), style="List Bullet")
        elif kind == "ol":
            for item in block[1]:
                doc.add_paragraph(strip_inline_md(item), style="List Number")
        elif kind == "table":
            rows = block[1]
            if not rows:
                continue
            ncols = max(len(r) for r in rows)
            table = doc.add_table(rows=len(rows), cols=ncols)
            table.style = "Light Grid Accent 1" if "Light Grid Accent 1" in [s.name for s in doc.styles] else table.style
            for r, row in enumerate(rows):
                for c in range(ncols):
                    cell_text = row[c] if c < len(row) else ""
                    table.cell(r, c).text = strip_inline_md(cell_text)

    doc.add_paragraph()
    footer_p = doc.add_paragraph(FOOTER_TEXT)
    for run in footer_p.runs:
        run.font.size = Pt(8)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))


def cmd_md2docx(args: argparse.Namespace) -> int:
    in_path = Path(args.input)
    out_path = Path(args.output)
    text = read_text(in_path)
    meta, body = split_front_matter(text)
    blocks = parse_markdown_blocks(body)
    try:
        build_docx_direct(meta, blocks, out_path)
        print(f"wrote {out_path} (direct python-docx)")
        return 0
    except Exception as exc:  # fallback to soffice
        print(f"direct python-docx build failed ({exc}); falling back to soffice html->docx", file=sys.stderr)
        with tempfile.TemporaryDirectory() as td:
            tmp_html = Path(td) / (in_path.stem + ".html")
            cmd_md2html(argparse.Namespace(input=str(in_path), output=str(tmp_html), lang=None, css=True))
            try:
                produced = soffice_convert(tmp_html, Path(td), "docx:MS Word 2007 XML")
            except Exception:
                odt = soffice_convert(tmp_html, Path(td), "odt")
                produced = soffice_convert(odt, Path(td), "docx:MS Word 2007 XML")
            out_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy(produced, out_path)
        print(f"wrote {out_path} (soffice fallback)")
        return 0


# --------------------------------------------------------------------------
# md2pptx
# --------------------------------------------------------------------------

def build_pptx(meta: dict, blocks, out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    title = meta.get("title", out_path.stem)
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        subtitle_bits = [v for k, v in meta.items() if k in ("author", "date", "lang") and v]
        slide.placeholders[1].text = " · ".join(subtitle_bits)

    current_slide = None
    current_body_lines: list[str] = []

    def flush():
        nonlocal current_slide, current_body_lines
        if current_slide is not None:
            tf = current_slide.placeholders[1].text_frame
            tf.word_wrap = True
            if current_body_lines:
                tf.text = current_body_lines[0]
                for line in current_body_lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
        current_body_lines = []

    for block in blocks:
        kind = block[0]
        if kind == "h" and block[1] in (1, 2):
            flush()
            current_slide = prs.slides.add_slide(content_layout)
            current_slide.shapes.title.text = strip_inline_md(block[2])
            current_body_lines = []
        elif current_slide is None:
            continue  # skip content before first heading (front matter already used)
        elif kind == "p":
            current_body_lines.append(strip_inline_md(block[1]))
        elif kind in ("ul", "ol"):
            current_body_lines.extend(strip_inline_md(x) for x in block[1])
        elif kind == "table":
            for row in block[1]:
                current_body_lines.append(" | ".join(strip_inline_md(c) for c in row))
        elif kind == "quote":
            current_body_lines.append(f"\u201c{strip_inline_md(block[1])}\u201d")
        elif kind == "code":
            current_body_lines.append(block[1])
        elif kind == "h":  # h3/h4 inside a slide -> sub-bullet label
            current_body_lines.append(strip_inline_md(block[2]).upper())
    flush()

    # Closing "Blackbox Note" reminder slide
    closing = prs.slides.add_slide(content_layout)
    closing.shapes.title.text = "Blackbox Note"
    tf = closing.placeholders[1].text_frame
    tf.word_wrap = True
    tf.text = "Review the Blackbox Note before acting on this deck's claims."
    p = tf.add_paragraph()
    p.text = "Tier discipline: content here is a readout, not settled truth, unless the source explicitly states a higher tier."
    p2 = tf.add_paragraph()
    p2.text = FOOTER_TEXT
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def cmd_md2pptx(args: argparse.Namespace) -> int:
    in_path = Path(args.input)
    out_path = Path(args.output)
    text = read_text(in_path)
    meta, body = split_front_matter(text)
    blocks = parse_markdown_blocks(body)
    build_pptx(meta, blocks, out_path)
    print(f"wrote {out_path}")
    return 0


# --------------------------------------------------------------------------
# tex2pdf
# --------------------------------------------------------------------------

def cmd_tex2pdf(args: argparse.Namespace) -> int:
    tex_path = Path(args.input).resolve()
    if not tex_path.exists():
        print(f"error: {tex_path} not found", file=sys.stderr)
        return 1
    src_text = read_text(tex_path)
    if has_thai(src_text):
        print("WARNING: source .tex contains Thai characters. pdflatex (the only working "
              "LaTeX engine in this environment) is EN-only and cannot render Thai glyphs — "
              "xelatex is missing and lualatex's fontspec/luaotfload is broken here. "
              "Thai LaTeX output is NOT supported by this tool; use md2pdf (soffice route) "
              "for Thai documents instead.", file=sys.stderr)

    pdflatex = shutil.which("pdflatex")
    if not pdflatex:
        print("error: pdflatex not found on PATH", file=sys.stderr)
        return 1

    work_dir = tex_path.parent
    base = tex_path.stem

    def run_pdflatex():
        return subprocess.run(
            [pdflatex, "-interaction=nonstopmode", "-halt-on-error", tex_path.name],
            cwd=str(work_dir), capture_output=True, text=True, timeout=180,
        )

    r1 = run_pdflatex()
    if r1.returncode != 0:
        print(r1.stdout[-4000:], file=sys.stderr)
        return r1.returncode

    bib_file = work_dir / f"{base}.bib"
    aux_file = work_dir / f"{base}.aux"
    bibtex = shutil.which("bibtex")
    if bib_file.exists() and bibtex and aux_file.exists():
        rb = subprocess.run([bibtex, base], cwd=str(work_dir), capture_output=True, text=True, timeout=120)
        if rb.returncode != 0:
            print("warning: bibtex reported errors:\n" + rb.stdout[-2000:], file=sys.stderr)

    r2 = run_pdflatex()
    if r2.returncode != 0:
        print(r2.stdout[-4000:], file=sys.stderr)
        return r2.returncode

    pdf_out = work_dir / f"{base}.pdf"
    print(f"wrote {pdf_out}")
    return 0


# --------------------------------------------------------------------------
# check (toolchain check)
# --------------------------------------------------------------------------

def _tool_version(cmd: list[str]) -> str:
    try:
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=15)
        out = (r.stdout or r.stderr).strip().splitlines()
        return out[0][:60] if out else "unknown"
    except Exception:
        return "unknown"


def cmd_check(args: argparse.Namespace) -> int:
    rows = []

    def add(name, present, version, needed_for, degrade):
        rows.append((name, "yes" if present else "NO", version if present else "-", needed_for, degrade))

    py_ok = sys.version_info >= (3, 10)
    add("python3", py_ok, ".".join(map(str, sys.version_info[:3])), "core: everything", "nothing runs")

    for mod, needed_for, degrade in [
        ("markdown", "md2html/md2pdf/md2docx/md2pptx source parsing", "no rendering at all"),
        ("docx", "md2docx direct build", "falls back to soffice html->docx"),
        ("pptx", "md2pptx", "pptx route unavailable"),
    ]:
        try:
            m = __import__(mod)
            add(mod, True, getattr(m, "__version__", "n/a"), needed_for, degrade)
        except Exception:
            add(mod, False, "-", needed_for, degrade)

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    add("soffice/libreoffice", bool(soffice), _tool_version([soffice, "--version"]) if soffice else "-",
        "md2pdf, md2docx fallback, Thai-capable PDF/DOCX", "no PDF route, no Thai rendering")

    pdflatex = shutil.which("pdflatex")
    add("pdflatex", bool(pdflatex), _tool_version([pdflatex, "--version"]) if pdflatex else "-",
        "tex2pdf (EN-only)", "no LaTeX route at all")

    xelatex = shutil.which("xelatex")
    add("xelatex", bool(xelatex), _tool_version([xelatex, "--version"]) if xelatex else "-",
        "Thai-capable LaTeX (not currently reachable)", "Thai LaTeX unavailable; use md2pdf instead")

    lualatex = shutil.which("lualatex")
    add("lualatex", bool(lualatex), "present but fontspec/luaotfload broken for Thai (verified 2026-09-04)",
        "Thai-capable LaTeX (currently broken)", "treat as unavailable for Thai; EN-only if used at all")

    pandoc = shutil.which("pandoc")
    add("pandoc", bool(pandoc), _tool_version([pandoc, "--version"]) if pandoc else "-",
        "optional alternate md conversion route", "not used by this tool; direct routes cover it")

    bibtex = shutil.which("bibtex")
    add("bibtex", bool(bibtex), _tool_version([bibtex, "--version"]) if bibtex else "-",
        "tex2pdf bibliography step", "tex2pdf skips bibliography if absent")

    gh = shutil.which("gh")
    add("gh", bool(gh), _tool_version([gh, "--version"]) if gh else "-",
        "release/archive workflow (not this tool)", "manual GitHub steps")

    curl = shutil.which("curl")
    add("curl", bool(curl), _tool_version([curl, "--version"]) if curl else "-",
        "citation checking scripts (not this tool)", "citation checks unavailable")

    thai_fonts = _tool_version(["fc-list"]) if shutil.which("fc-list") else ""
    has_noto_thai = False
    if shutil.which("fc-list"):
        r = subprocess.run(["fc-list"], capture_output=True, text=True, timeout=15)
        has_noto_thai = "Noto Sans Thai" in r.stdout or "Noto Serif Thai" in r.stdout
    add("Noto Thai fonts", has_noto_thai, "Noto Sans/Serif/Looped Thai" if has_noto_thai else "-",
        "Thai rendering via soffice route", "Thai text may render with fallback/missing glyphs")

    core_missing = [r[0] for r in rows if r[0] in ("python3", "markdown") and r[1] == "NO"]

    col_w = [max(len(r[i]) for r in rows + [("tool", "present", "version", "needed-for", "degrade")]) for i in range(5)]
    header = ("tool", "present", "version", "needed-for", "degrade")
    def fmt(row):
        return " | ".join(str(c).ljust(col_w[i]) for i, c in enumerate(row))
    print(fmt(header))
    print("-+-".join("-" * w for w in col_w))
    for row in rows:
        print(fmt(row))

    if core_missing:
        print(f"\nCORE tool(s) missing: {', '.join(core_missing)} — see TOOLCHAIN.md", file=sys.stderr)
        return 1
    return 0


# --------------------------------------------------------------------------
# all
# --------------------------------------------------------------------------

def cmd_all(args: argparse.Namespace) -> int:
    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)
    rc = 0
    for in_str in args.inputs:
        in_path = Path(in_str)
        if not in_path.exists():
            print(f"skip (not found): {in_path}", file=sys.stderr)
            continue
        stem = in_path.stem
        for suffix, fn in (("pdf", cmd_md2pdf), ("docx", cmd_md2docx), ("pptx", cmd_md2pptx)):
            out_path = out_dir / f"{stem}.{suffix}"
            try:
                fn(argparse.Namespace(input=str(in_path), output=str(out_path)))
            except Exception as exc:
                print(f"FAILED {suffix} for {in_path}: {exc}", file=sys.stderr)
                rc = 1
        # copy any sibling LaTeX PDF of the same stem if present
        tex_pdf_candidates = list(in_path.parent.glob(f"{stem}*.pdf"))
        for cand in tex_pdf_candidates:
            dest = out_dir / f"{stem}_latex_{cand.name}"
            if not dest.exists():
                shutil.copy(cand, dest)
                print(f"copied latex pdf {cand} -> {dest}")
    return rc


# --------------------------------------------------------------------------
# deliver-manifest
# --------------------------------------------------------------------------

CHANNEL_MAP = {
    ".pdf": "google_drive",
    ".docx": "google_drive",
    ".pptx": "canva_or_pptx_attachment",
    ".html": "google_drive_or_email_body",
}


def cmd_deliver_manifest(args: argparse.Namespace) -> int:
    scan_dir = Path(args.dir)
    out_path = Path(args.output)
    entries = []
    if scan_dir.exists():
        for p in sorted(scan_dir.rglob("*")):
            if p.is_file() and p.suffix.lower() in CHANNEL_MAP:
                entries.append({
                    "file": str(p),
                    "suggested_channel": CHANNEL_MAP[p.suffix.lower()],
                    "requires_human_approval_before_send": True,
                    "k_state_note": "Include the K-state / tier disclaimer line in the delivery message body before sharing.",
                })
    manifest = {
        "render_tool_version": VERSION,
        "generated_from": str(scan_dir),
        "files": entries,
        "note": "This manifest lists rendered files and a suggested delivery channel per file. "
                "It does not send or share anything itself — an AI session with the matching "
                "connector reads this and performs delivery, gated by human approval per "
                "plugins/glosa/skills/glosa-deliver/SKILL.md.",
    }
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"wrote {out_path} ({len(entries)} file(s))")
    return 0


# --------------------------------------------------------------------------
# CLI wiring
# --------------------------------------------------------------------------

def main(argv=None) -> int:
    parser = argparse.ArgumentParser(prog="render.py", description="glosa document-rendering layer")
    sub = parser.add_subparsers(dest="command", required=True)

    p = sub.add_parser("md2html")
    p.add_argument("input")
    p.add_argument("output")
    p.add_argument("--lang", choices=["th", "en"], default=None)
    p.add_argument("--css", action="store_true")
    p.set_defaults(func=cmd_md2html)

    p = sub.add_parser("md2pdf")
    p.add_argument("input")
    p.add_argument("output")
    p.set_defaults(func=cmd_md2pdf)

    p = sub.add_parser("md2docx")
    p.add_argument("input")
    p.add_argument("output")
    p.set_defaults(func=cmd_md2docx)

    p = sub.add_parser("md2pptx")
    p.add_argument("input")
    p.add_argument("output")
    p.set_defaults(func=cmd_md2pptx)

    p = sub.add_parser("tex2pdf")
    p.add_argument("input")
    p.set_defaults(func=cmd_tex2pdf)

    p = sub.add_parser("all")
    p.add_argument("inputs", nargs="+")
    p.add_argument("--out", default="dist/")
    p.set_defaults(func=cmd_all)

    p = sub.add_parser("check")
    p.set_defaults(func=cmd_check)

    p = sub.add_parser("deliver-manifest")
    p.add_argument("output")
    p.add_argument("--dir", default="dist")
    p.set_defaults(func=cmd_deliver_manifest)

    args = parser.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
